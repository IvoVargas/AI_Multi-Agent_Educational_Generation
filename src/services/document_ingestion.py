from __future__ import annotations

import json
import mimetypes
import re
from pathlib import Path
from typing import Any, Dict, List
from uuid import uuid4

from src.utils.logging_utils import get_logger

logger = get_logger(__name__)

TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".csv", ".json", ".yaml", ".yml", ".rst"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".potx"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | IMAGE_EXTENSIONS | DOCUMENT_EXTENSIONS


def _normalize_whitespace(text: str) -> str:
    return re.sub(r"\n{3,}", "\n\n", re.sub(r"[ \t]+", " ", text)).strip()


def _extract_plain_text(path: Path) -> str:
    try:
        return _normalize_whitespace(path.read_text(encoding="utf-8"))
    except UnicodeDecodeError:
        return _normalize_whitespace(path.read_text(encoding="latin-1", errors="ignore"))


def _extract_json_text(path: Path) -> str:
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        return _normalize_whitespace(json.dumps(data, ensure_ascii=False, indent=2))
    except Exception:
        return _extract_plain_text(path)


def _extract_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except Exception:
        logger.warning("pypdf not available. Skipping PDF text extraction for %s", path.name)
        return ""

    chunks: List[str] = []
    try:
        reader = PdfReader(str(path))
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(text)
    except Exception:
        logger.exception("Failed to extract PDF text from %s", path.name)
        return ""
    return _normalize_whitespace("\n\n".join(chunks))


def _extract_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except Exception:
        logger.warning("python-docx not available. Skipping DOCX extraction for %s", path.name)
        return ""

    try:
        doc = Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if values:
                    parts.append(" | ".join(values))
        return _normalize_whitespace("\n".join(parts))
    except Exception:
        logger.exception("Failed to extract DOCX text from %s", path.name)
        return ""


def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception:
        logger.warning("python-pptx not available. Skipping PPTX extraction for %s", path.name)
        return ""

    try:
        prs = Presentation(str(path))
        slides: List[str] = []
        for index, slide in enumerate(prs.slides, start=1):
            parts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text).strip():
                    parts.append(str(shape.text).strip())
            if parts:
                slides.append(f"Slide {index}\n" + "\n".join(parts))
        return _normalize_whitespace("\n\n".join(slides))
    except Exception:
        logger.exception("Failed to extract PPTX text from %s", path.name)
        return ""


def extract_text_from_file(file_path: str) -> str:
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext in {".json"}:
        return _extract_json_text(path)
    if ext in TEXT_EXTENSIONS:
        return _extract_plain_text(path)
    if ext == ".pdf":
        return _extract_pdf_text(path)
    if ext == ".docx":
        return _extract_docx_text(path)
    if ext in {".pptx", ".potx"}:
        return _extract_pptx_text(path)
    return ""


def make_chunks(text: str, chunk_size: int = 1200, overlap: int = 180) -> List[str]:
    cleaned = _normalize_whitespace(text)
    if not cleaned:
        return []

    paragraphs = [p.strip() for p in re.split(r"\n\n+", cleaned) if p.strip()]
    chunks: List[str] = []
    current = ""
    for paragraph in paragraphs:
        candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if len(candidate) <= chunk_size:
            current = candidate
            continue
        if current:
            chunks.append(current)
        if len(paragraph) <= chunk_size:
            current = paragraph
        else:
            start = 0
            while start < len(paragraph):
                end = min(len(paragraph), start + chunk_size)
                piece = paragraph[start:end].strip()
                if piece:
                    chunks.append(piece)
                start = max(end - overlap, start + 1)
            current = ""
    if current:
        chunks.append(current)
    return chunks


def ingest_file(file_path: str) -> Dict[str, Any]:
    path = Path(file_path)
    ext = path.suffix.lower()
    mime_type, _ = mimetypes.guess_type(str(path))
    text = extract_text_from_file(str(path)) if ext in TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS else ""
    chunks = make_chunks(text) if text else []

    return {
        "id": f"file_{uuid4().hex[:8]}",
        "name": path.name,
        "path": str(path),
        "extension": ext,
        "mime_type": mime_type or "application/octet-stream",
        "text": text,
        "text_excerpt": text[:1500],
        "chunk_count": len(chunks),
        "chunks": chunks,
        "is_visual": ext in IMAGE_EXTENSIONS,
        "is_template": ext in {".pptx", ".potx"},
        "is_textual": bool(text.strip()),
        "status": "ingested",
    }


def score_chunk_relevance(chunk: str, query: str) -> int:
    if not chunk.strip() or not query.strip():
        return 0
    tokens = {token for token in re.findall(r"\w+", query.lower()) if len(token) > 2}
    if not tokens:
        return 0
    lower_chunk = chunk.lower()
    score = sum(1 for token in tokens if token in lower_chunk)
    if lower_chunk[:120].count("slide"):
        score += 1
    return score


def select_relevant_chunks(knowledge_chunks: List[Dict[str, Any]], query: str, top_n: int = 5) -> List[Dict[str, Any]]:
    scored: List[Dict[str, Any]] = []
    for item in knowledge_chunks or []:
        score = score_chunk_relevance(str(item.get("text", "")), query)
        if score > 0:
            scored.append({**item, "score": score})
    scored.sort(key=lambda item: (item.get("score", 0), len(str(item.get("text", "")))), reverse=True)
    return scored[:top_n]


def supported_file_message(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext in SUPPORTED_EXTENSIONS:
        return ""
    return f"O ficheiro '{Path(file_path).name}' não é suportado nesta versão. Usa TXT, MD, PDF, DOCX, PPTX ou imagens."
