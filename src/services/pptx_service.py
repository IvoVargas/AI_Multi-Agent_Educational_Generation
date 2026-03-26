from __future__ import annotations

from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import Any, Dict, List

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PRIMARY = RGBColor(33, 37, 41)
ACCENT = RGBColor(52, 120, 246)
MUTED = RGBColor(108, 117, 125)
LIGHT_BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)

SLIDE_WIDTH_IN = 10.0
SLIDE_HEIGHT_IN = 7.5


def _set_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def _add_top_bar(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_WIDTH_IN),
        Inches(0.28),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def _add_footer(slide, text: str = "Gerado automaticamente pelo protótipo") -> None:
    textbox = slide.shapes.add_textbox(
        Inches(0.35),
        Inches(7.02),
        Inches(9.2),
        Inches(0.22),
    )
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def _add_textbox(slide, left, top, width, height, text=""):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = text
    return textbox


def _add_title_text(slide, title: str, size: int = 24) -> None:
    textbox = _add_textbox(
        slide,
        Inches(0.6),
        Inches(0.55),
        Inches(8.7),
        Inches(0.7),
        title,
    )
    p = textbox.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = PRIMARY


def _add_subtitle_text(slide, subtitle: str) -> None:
    textbox = _add_textbox(
        slide,
        Inches(0.75),
        Inches(2.2),
        Inches(8.4),
        Inches(0.8),
        subtitle,
    )
    p = textbox.text_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER


def _crop_image_to_ratio(image_path: str, target_ratio: float) -> str:
    with Image.open(image_path) as img:
        img = img.convert("RGB")
        width, height = img.size
        current_ratio = width / height

        if current_ratio > target_ratio:
            new_width = int(height * target_ratio)
            left = (width - new_width) // 2
            img = img.crop((left, 0, left + new_width, height))
        else:
            new_height = int(width / target_ratio)
            top = (height - new_height) // 2
            img = img.crop((0, top, width, top + new_height))

        temp_file = NamedTemporaryFile(delete=False, suffix=".png")
        img.save(temp_file.name, format="PNG")
        return temp_file.name


def _add_image_or_visual_text(
    slide,
    left,
    top,
    width,
    height,
    image_path: str | None,
    visual_description: str,
    title: str = "Apoio visual",
) -> None:
    container = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    container.fill.solid()
    container.fill.fore_color.rgb = LIGHT_BG
    container.line.color.rgb = ACCENT

    title_box = slide.shapes.add_textbox(
        left + Inches(0.12),
        top + Inches(0.08),
        width - Inches(0.24),
        Inches(0.28),
    )
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.bold = True
    tp.font.size = Pt(14)
    tp.font.color.rgb = ACCENT
    tp.alignment = PP_ALIGN.CENTER

    image_area_left = left + Inches(0.12)
    image_area_top = top + Inches(0.45)
    image_area_width = width - Inches(0.24)
    image_area_height = height - Inches(0.95)

    if image_path and Path(image_path).exists():
        target_ratio = image_area_width / image_area_height
        cropped_path = _crop_image_to_ratio(image_path, target_ratio)
        slide.shapes.add_picture(
            cropped_path,
            image_area_left,
            image_area_top,
            width=image_area_width,
            height=image_area_height,
        )
    else:
        text_box = slide.shapes.add_textbox(
            image_area_left,
            image_area_top,
            image_area_width,
            image_area_height,
        )
        tf = text_box.text_frame
        tf.clear()
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = visual_description or "Sugestão visual não disponível."
        p.font.size = Pt(12)
        p.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER

    caption_box = slide.shapes.add_textbox(
        left + Inches(0.12),
        top + height - Inches(0.35),
        width - Inches(0.24),
        Inches(0.20),
    )
    cp = caption_box.text_frame.paragraphs[0]
    cp.text = "Imagem gerada automaticamente" if image_path and Path(image_path).exists() else "Descrição visual"
    cp.font.size = Pt(9)
    cp.font.color.rgb = MUTED
    cp.alignment = PP_ALIGN.CENTER


def _add_body_box(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.55),
        Inches(1.45),
        Inches(5.55),
        Inches(5.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 225, 230)


def _truncate_bullet(text: str, max_chars: int = 110) -> str:
    text = (text or "").strip()
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


def _add_bullets(slide, bullets: List[str]) -> None:
    textbox = _add_textbox(
        slide,
        Inches(0.8),
        Inches(1.8),
        Inches(4.95),
        Inches(4.45),
        "",
    )

    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True

    for idx, bullet in enumerate((bullets or ["Conteúdo não disponível"])[:5]):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = _truncate_bullet(bullet, 110)
        p.font.size = Pt(18)
        p.font.color.rgb = PRIMARY
        p.level = 0
        p.space_after = Pt(8)


def _add_section_label(slide, text: str) -> None:
    textbox = _add_textbox(
        slide,
        Inches(0.65),
        Inches(1.28),
        Inches(2.2),
        Inches(0.25),
        text,
    )
    p = textbox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACCENT


def _build_title_slide(slide, item: Dict[str, Any], metadata: Dict[str, Any]) -> None:
    _set_slide_background(slide)
    _add_top_bar(slide)

    _add_title_text(slide, item.get("title", "Apresentação"), size=28)
    _add_subtitle_text(
        slide,
        metadata.get("presentation_goal", "") or "Apresentação educativa",
    )

    _add_image_or_visual_text(
        slide,
        Inches(1.2),
        Inches(3.1),
        Inches(7.6),
        Inches(2.3),
        item.get("image_path"),
        item.get("visual_description", ""),
        title="Imagem de capa",
    )

    audience = metadata.get("target_audience", "").strip()
    level = metadata.get("education_level", "").strip()
    info = " | ".join(x for x in [audience, level] if x)

    if info:
        info_box = _add_textbox(
            slide,
            Inches(0.9),
            Inches(6.0),
            Inches(8.2),
            Inches(0.3),
            info,
        )
        p = info_box.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.CENTER

    _add_footer(slide)


def _build_content_slide(slide, item: Dict[str, Any]) -> None:
    _set_slide_background(slide)
    _add_top_bar(slide)

    _add_title_text(slide, item.get("title", "Slide"), size=23)
    _add_body_box(slide)
    _add_bullets(slide, item.get("bullets", []))

    _add_image_or_visual_text(
        slide,
        Inches(6.35),
        Inches(1.45),
        Inches(3.1),
        Inches(5.2),
        item.get("image_path"),
        item.get("visual_description", ""),
        title="Apoio visual",
    )

    _add_footer(slide)


def build_presentation(
    slide_plan: List[Dict[str, Any]],
    metadata: Dict[str, Any],
    output_path: str,
) -> str:
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    blank_layout = prs.slide_layouts[6]

    for item in slide_plan:
        slide = prs.slides.add_slide(blank_layout)
        kind = item.get("kind", "content")

        if kind == "title":
            _build_title_slide(slide, item, metadata)
        else:
            _build_content_slide(slide, item)

    prs.save(str(output_file))
    return str(output_file)